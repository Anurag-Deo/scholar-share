# app/services/html_export_service.py

from pathlib import Path
import asyncio
from bs4 import BeautifulSoup
from pptx import Presentation
from playwright.async_api import async_playwright


class HTMLExportService:
    """
    A service to render HTML slides into PDF and PPTX formats using async APIs.
    """

    async def html_to_pdf(self, html_path: Path, pdf_path: Path) -> Path:
        """Converts an HTML file to a PDF using a headless browser, asynchronously."""
        async with async_playwright() as pw:
            browser = await pw.chromium.launch()
            page = await browser.new_page()

            # --- FIX: Resolve the relative path to an absolute one before creating a URI ---
            absolute_html_path_uri = html_path.resolve().as_uri()
            await page.goto(absolute_html_path_uri, wait_until="networkidle")

            await page.emulate_media(media="screen")
            await page.pdf(path=str(pdf_path), format="A4", print_background=True)
            await browser.close()
        return pdf_path

    def _sync_html_to_ppt(self, html_path: Path, ppt_path: Path) -> Path:
        """
        Internal synchronous method to parse HTML and build a .pptx file.
        This contains blocking I/O and should not be called directly in an async context.
        """
        prs = Presentation()
        with open(html_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")

        for section in soup.select("section"):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            title_element = section.find(["h1", "h2", "h3"])
            if title_element:
                slide.shapes.title.text = title_element.get_text(strip=True)

            if len(slide.placeholders) > 1:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                for bullet in section.find_all("li"):
                    p = tf.add_paragraph()
                    p.text = bullet.get_text(strip=True)
                    p.level = 0

        prs.save(ppt_path)
        return ppt_path

    async def html_to_ppt(self, html_path: Path, ppt_path: Path) -> Path:
        """
        Asynchronously triggers the synchronous conversion of HTML to PPTX
        by running it in a separate thread to avoid blocking the event loop.
        """
        return await asyncio.to_thread(self._sync_html_to_ppt, html_path, ppt_path)


html_export_service = HTMLExportService()
